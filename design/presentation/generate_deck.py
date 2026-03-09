#!/usr/bin/env python3
"""
Generate eGramSabha Hackathon Submission Deck (33 slides).
Opens the template PPTX (14 slides with backgrounds), fills template
slides 1-13, leaves template slide 14 untouched, appends extra slides
(including our title slide and closing slide), then reorders via XML.

Final order:
  1:  Template landing (fields filled)
  2:  Our eGramSabha title (extra)
  3-14: Template slides 2-13
  15-31: Extra content slides
  32: Our Thank You (extra)
  33: Template ending (untouched)

Usage:
    python generate_deck.py
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from styles import (
    DARK_TEXT, WHITE, GREEN_PRIMARY, GREEN_LIGHT, ORANGE_ACCENT,
    BLUE_AWS, BLUE_LIGHT, GRAY_MED, RED_URGENT, YELLOW_READY,
    TABLE_HEADER_BG, TABLE_ALT_BG, GRAY_LIGHT,
    PLACEHOLDER_BG, PLACEHOLDER_BORDER,
    HEADER_BAR_COLOR, HEADER_BAR_HEIGHT,
    FONT_TITLE, FONT_BODY, FONT_MONO,
    SIZE_SLIDE_TITLE, SIZE_SUBTITLE, SIZE_BODY, SIZE_BODY_SMALL,
    SIZE_TABLE, SIZE_TABLE_HEADER, SIZE_CAPTION, SIZE_BIG_STAT,
    CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT,
    SLIDE_WIDTH, SLIDE_HEIGHT,
    set_font, add_titled_textbox, add_bullet_point, add_paragraph,
    style_table, set_cell_text,
)
from content import (
    SLIDE1, SLIDE_TITLE, SLIDE2, SLIDE3, SLIDE4, SLIDE5, SLIDE6, SLIDE7, SLIDE8,
    SLIDE9, SLIDE10, SLIDE11, SLIDE12, SLIDE13, SLIDE13_LINKS, SLIDE14,
    SLIDE15, SLIDE16, SLIDE17, SLIDE18, SLIDE19, SLIDE20,
    SLIDE21, SLIDE22, SLIDE23, SLIDE24, SLIDE25, SLIDE26,
    SLIDE27, GALLERY_SLIDES, SLIDE30, SLIDE31,
)

# Paths
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.abspath(os.path.join(SCRIPT_DIR, "..", ".."))
TEMPLATE_PATH = os.path.join(PROJECT_ROOT, "design",
                             "Prototype Development Submission _ AWS AI for Bharat Hackathon.pptx")
OUTPUT_PATH = os.path.join(PROJECT_ROOT, "design", "eGramSabha_Hackathon_Submission.pptx")
LOGO_PATH = os.path.join(SCRIPT_DIR, "assets", "logo.png")
SCREENSHOTS_DIR = os.path.join(SCRIPT_DIR, "screenshots")
DIAGRAMS_DIR = os.path.join(SCRIPT_DIR, "diagrams")
TEMPLATE_BG_PATH = os.path.join(SCRIPT_DIR, "assets", "template_bg.png")


# ============================================================
# HELPERS
# ============================================================

def diagram_path(name):
    p = os.path.join(DIAGRAMS_DIR, f"{name}.png")
    return p if os.path.exists(p) else None


def get_image_size(path):
    """Return (width, height) in pixels for an image file."""
    from PIL import Image
    with Image.open(path) as img:
        return img.size


def fit_image(img_path, max_w, max_h):
    """Calculate image dimensions that fit within max_w x max_h preserving aspect ratio.
    Returns (width, height) as Emu values."""
    pw, ph = get_image_size(img_path)
    ratio = pw / ph
    # Try fitting to width first
    w = max_w
    h = int(w / ratio)
    if h > max_h:
        # Too tall, fit to height instead
        h = max_h
        w = int(h * ratio)
    return w, h


def screenshot_path(name):
    for ext in ('.png', '.jpg', '.jpeg'):
        p = os.path.join(SCREENSHOTS_DIR, f"{name}{ext}")
        if os.path.exists(p):
            return p
    return None


def clear_slide_text(slide):
    """Remove all text boxes (shape_type 17) from a slide, keep background picture."""
    for shape in [s for s in slide.shapes if s.shape_type == 17]:
        sp = shape._element
        sp.getparent().remove(sp)


def add_textbox(slide, left, top, width, height, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    return tf


def add_table(slide, rows, cols, left, top, width, height):
    return slide.shapes.add_table(rows, cols, left, top, width, height).table


def add_flow_box(slide, left, top, width, height, text,
                 fill_color=GREEN_PRIMARY, text_color=WHITE, font_size=Pt(9)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    set_font(run, size=font_size, bold=True, color=text_color)
    return shape


def add_arrow(slide, left, top, width, height=Inches(0.2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY_MED
    shape.line.fill.background()
    return shape


def add_placeholder(slide, left, top, width, height, label):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PLACEHOLDER_BG
    shape.line.color.rgb = PLACEHOLDER_BORDER
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = label
    set_font(run, SIZE_BODY_SMALL, False, GRAY_MED)
    return shape


def add_callout_box(slide, left, top, width, height, text,
                    fill_color=GREEN_PRIMARY, font_size=Pt(10)):
    """Rounded rectangle with centered white bold text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    set_font(run, font_size, True, WHITE)
    return shape


def add_section_header(tf_or_slide, text, left=None, top=None, width=None,
                       color=GREEN_PRIMARY, size=Pt(12)):
    """Add a section header. If left/top/width provided, creates textbox on slide."""
    if left is not None:
        tf = add_textbox(tf_or_slide, left, top, width, Inches(0.3))
        p = tf.paragraphs[0]
    else:
        p = tf_or_slide.paragraphs[0] if not tf_or_slide.paragraphs[0].text else tf_or_slide.add_paragraph()
    run = p.add_run()
    run.text = text
    set_font(run, size, True, color)
    return p


def add_extra_slide(prs, title_text=""):
    """Add slide beyond template: uses template background image + title."""
    slide = prs.slides.add_slide(prs.slide_layouts[10])  # BLANK
    # Use the same background image as template slides
    if os.path.exists(TEMPLATE_BG_PATH):
        slide.shapes.add_picture(TEMPLATE_BG_PATH, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    else:
        # Fallback: white bg + green header bar
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.fill.background()
        hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, HEADER_BAR_HEIGHT)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = HEADER_BAR_COLOR
        hdr.line.fill.background()
    # Title on slide (positioned to match template title placement)
    if title_text:
        tf = add_textbox(slide, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Inches(0.5))
        run = tf.paragraphs[0].add_run()
        run.text = title_text
        set_font(run, SIZE_SLIDE_TITLE, True, DARK_TEXT)
    return slide


def _title_on_template(slide, title_text, color=DARK_TEXT):
    """Standard title on a template slide (below green header band)."""
    tf = add_textbox(slide, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Inches(0.5))
    run = tf.paragraphs[0].add_run()
    run.text = title_text
    set_font(run, SIZE_SLIDE_TITLE, True, color)


def _feature_list_with_screenshot(slide, title, features, screenshot_label, screenshot_name=None):
    """Common pattern for portal slides: features left with colored indicators, screenshot right."""
    clear_slide_text(slide)
    _title_on_template(slide, title)

    feature_colors = [GREEN_PRIMARY, BLUE_LIGHT, ORANGE_ACCENT,
                      RGBColor(0x7B, 0x1F, 0xA2), RED_URGENT, GREEN_PRIMARY]

    y = Inches(1.4)
    for i, (name, desc) in enumerate(features):
        color = feature_colors[i % len(feature_colors)]

        # Colored circle indicator
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(0.3), y + Inches(0.03),
                                     Inches(0.14), Inches(0.14))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()

        # Feature name (slightly larger)
        tf_name = add_textbox(slide, Inches(0.5), y - Inches(0.02), Inches(4.9), Inches(0.22))
        run = tf_name.paragraphs[0].add_run()
        run.text = name
        set_font(run, Pt(10), True, color)

        # Thin colored accent line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(0.5), y + Inches(0.2),
                                      Inches(4.9), Pt(1.5))
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()

        # Description below the line
        tf_desc = add_textbox(slide, Inches(0.5), y + Inches(0.22), Inches(4.9), Inches(0.25))
        run2 = tf_desc.paragraphs[0].add_run()
        run2.text = desc
        set_font(run2, Pt(7.5), color=DARK_TEXT)

        y += Inches(0.5)

    # Try real screenshot, fall back to placeholder
    sp = screenshot_path(screenshot_name) if screenshot_name else None
    img_area_x = Inches(5.7)
    img_area_y = Inches(1.4)
    img_area_w = Inches(4.0)
    img_area_h = Inches(3.6)
    if sp:
        try:
            w, h = fit_image(sp, img_area_w, img_area_h)
            # Center within the area
            x_off = img_area_x + (img_area_w - w) // 2
            y_off = img_area_y + (img_area_h - h) // 2
            slide.shapes.add_picture(sp, x_off, y_off, w, h)
            return
        except Exception:
            pass
    add_placeholder(slide, img_area_x, img_area_y, img_area_w, img_area_h,
                    f"[Screenshot]\n{screenshot_label}")


# ============================================================
# TEMPLATE SLIDE BUILDERS (1-13)
# ============================================================

def build_slide1(slide):
    """Template Landing (fields only)."""
    clear_slide_text(slide)
    d = SLIDE1

    # Only fill the 3 template text fields at their original positions.
    # Keep template background intact — no logo, tagline, or subtitle here.

    # Team Name at y~3.77"
    tf3 = add_textbox(slide, Emu(383675), Emu(3446998), Emu(8520600), Emu(443700))
    run = tf3.paragraphs[0].add_run()
    run.text = f"Team Name :  {d['team_name']}"
    set_font(run, Pt(16), True, DARK_TEXT)

    # Team Leader at y~4.22"
    tf4 = add_textbox(slide, Emu(383663), Emu(3854458), Emu(8520600), Emu(443700))
    run = tf4.paragraphs[0].add_run()
    run.text = f"Team Leader Name :  {d['team_leader']}"
    set_font(run, Pt(16), True, DARK_TEXT)

    # Problem Statement at y~4.70"
    tf5 = add_textbox(slide, Emu(383680), Emu(4298160), Emu(8520600), Emu(408600))
    run = tf5.paragraphs[0].add_run()
    run.text = f"Problem Statement :  {d['problem_statement']}"
    set_font(run, Pt(16), True, DARK_TEXT)


def build_our_title(prs):
    """Our eGramSabha Title Slide (extra)."""
    slide = add_extra_slide(prs, "")
    d = SLIDE_TITLE

    # Logo centered (below template header area)
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(4.0), Inches(0.6), Inches(1.5), Inches(1.5))

    # Project name
    tf = add_textbox(slide, Inches(0.5), Inches(2.2), Inches(9.0), Inches(0.7))
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = d["project_name"]
    set_font(run, Pt(36), True, GREEN_PRIMARY)

    # Tagline
    tf2 = add_textbox(slide, Inches(0.5), Inches(2.9), Inches(9.0), Inches(0.5))
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf2.paragraphs[0].add_run()
    run.text = d["tagline"]
    set_font(run, Pt(16), True, DARK_TEXT, italic=True)

    # Subtitle
    tf3 = add_textbox(slide, Inches(0.5), Inches(3.5), Inches(9.0), Inches(0.4))
    tf3.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf3.paragraphs[0].add_run()
    run.text = d["subtitle"]
    set_font(run, Pt(14), color=GRAY_MED)

    # Links row
    links = SLIDE13_LINKS
    link_items = []
    if links["github_repo"]:
        link_items.append(("GitHub", links["github_repo"]))
    if links["demo_video"]:
        link_items.append(("Demo Video", links["demo_video"]))
    if links["live_prototype"]:
        link_items.append(("Live Demo", links["live_prototype"]))

    tf_links = add_textbox(slide, Inches(0.5), Inches(3.85), Inches(9.0), Inches(0.9))
    for idx, (label, url) in enumerate(link_items):
        p = tf_links.paragraphs[0] if idx == 0 else tf_links.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        run_label = p.add_run()
        run_label.text = f"{label}: "
        set_font(run_label, Pt(9), bold=True, color=DARK_TEXT)
        run_url = p.add_run()
        run_url.text = url
        set_font(run_url, Pt(9), color=BLUE_LIGHT)
        run_url.hyperlink.address = url

    # Credentials
    cred_items = [
        ("Admin Login", "admin / AdminPassword123"),
        ("Official Login", "password same as configured username"),
    ]
    for label, value in cred_items:
        p = tf_links.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        run_label = p.add_run()
        run_label.text = f"{label}: "
        set_font(run_label, Pt(9), bold=True, color=DARK_TEXT)
        run_val = p.add_run()
        run_val.text = value
        set_font(run_val, Pt(9), color=GRAY_MED, name=FONT_MONO)

    # Event name
    tf4 = add_textbox(slide, Inches(0.5), Inches(4.85), Inches(9.0), Inches(0.4))
    tf4.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf4.paragraphs[0].add_run()
    run.text = d["event"]
    set_font(run, Pt(11), color=GRAY_MED, italic=True)

    return slide


def build_slide2(slide):
    """Problem Statement."""
    clear_slide_text(slide)
    d = SLIDE2
    _title_on_template(slide, d["title"])

    # Headline
    tf2 = add_textbox(slide, CONTENT_LEFT, Inches(1.4), CONTENT_WIDTH, Inches(0.55))
    run = tf2.paragraphs[0].add_run()
    run.text = d["headline"]
    set_font(run, Pt(9), color=DARK_TEXT, italic=True)

    # Stats row
    colors = [RED_URGENT, ORANGE_ACCENT, BLUE_LIGHT, GREEN_PRIMARY, GRAY_MED]
    for i, (stat, label) in enumerate(d["stats"]):
        x = Inches(0.4 + i * 1.9)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x, Inches(2.0), Inches(1.7), Inches(0.75))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i]
        shape.line.fill.background()
        tf_s = shape.text_frame
        tf_s.word_wrap = True
        tf_s.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf_s.paragraphs[0].add_run()
        run.text = stat
        set_font(run, Pt(18), True, WHITE)
        p2 = tf_s.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = label
        set_font(run2, Pt(7), color=WHITE)

    # Barriers table
    table = add_table(slide, len(d["barriers"]) + 1, 2,
                      Inches(0.4), Inches(2.9), Inches(9.2), Inches(2.3))
    set_cell_text(table, 0, 0, "Barrier", SIZE_TABLE_HEADER, True, WHITE)
    set_cell_text(table, 0, 1, "Impact", SIZE_TABLE_HEADER, True, WHITE)
    for i, (barrier, impact) in enumerate(d["barriers"]):
        set_cell_text(table, i + 1, 0, barrier)
        set_cell_text(table, i + 1, 1, impact)
    style_table(table)
    table.columns[0].width = Inches(4.0)
    table.columns[1].width = Inches(5.2)


def build_slide3(slide):
    """Solution Overview."""
    clear_slide_text(slide)
    d = SLIDE3
    _title_on_template(slide, d["title"])

    # Overview
    tf = add_textbox(slide, CONTENT_LEFT, Inches(1.4), CONTENT_WIDTH, Inches(0.7))
    run = tf.paragraphs[0].add_run()
    run.text = d["overview"]
    set_font(run, Pt(10), color=DARK_TEXT)

    # 3 Portal boxes
    portal_colors = [GREEN_PRIMARY, BLUE_LIGHT, ORANGE_ACCENT]
    for i, ((name, desc), color) in enumerate(zip(d["portals"], portal_colors)):
        x = Inches(0.3) + i * Inches(3.2)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x, Inches(2.15), Inches(3.0), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf_p = shape.text_frame
        tf_p.word_wrap = True
        tf_p.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf_p.paragraphs[0].add_run()
        run.text = name
        set_font(run, Pt(12), True, WHITE)
        p2 = tf_p.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = desc
        set_font(run2, Pt(7.5), color=WHITE)

    # Lifecycle diagram
    add_section_header(slide, "5-Step Lifecycle", Inches(0.3), Inches(3.5), Inches(4.0))
    wf = diagram_path("workflow")
    if wf:
        slide.shapes.add_picture(wf, Inches(0.3), Inches(3.85), Inches(9.3), Inches(0.7))
    else:
        tf4 = add_textbox(slide, CONTENT_LEFT, Inches(3.85), CONTENT_WIDTH, Inches(0.5))
        run = tf4.paragraphs[0].add_run()
        run.text = d["lifecycle"]
        set_font(run, Pt(10), True, GREEN_PRIMARY)


def build_slide4(slide):
    """5-Step Gram Sabha Lifecycle."""
    clear_slide_text(slide)
    d = SLIDE4
    _title_on_template(slide, d["title"])

    step_colors = [GREEN_PRIMARY, BLUE_LIGHT, ORANGE_ACCENT,
                   RGBColor(0x7B, 0x1F, 0xA2), RED_URGENT]

    # Horizontal card rows — each row: colored label on left + description on right
    label_w = Inches(2.2)
    desc_x = Inches(2.5)
    desc_w = Inches(7.0)
    row_h = Inches(0.62)
    y_start = Inches(1.4)

    for i, (step_name, keyword, desc) in enumerate(d["steps"]):
        y = y_start + i * (row_h + Inches(0.08))

        # Colored rounded rect with step name + keyword
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(0.3), y, label_w, row_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = step_colors[i]
        shape.line.fill.background()
        tf_s = shape.text_frame
        tf_s.word_wrap = True
        tf_s.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf_s.paragraphs[0].add_run()
        run.text = step_name
        set_font(run, Pt(10), True, WHITE)
        p2 = tf_s.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(0)
        run2 = p2.add_run()
        run2.text = keyword
        set_font(run2, Pt(8), color=WHITE)

        # Description text on right
        tf_d = add_textbox(slide, desc_x, y + Inches(0.05), desc_w, row_h - Inches(0.1))
        run = tf_d.paragraphs[0].add_run()
        run.text = desc
        set_font(run, Pt(8), color=DARK_TEXT)


def build_slide5(slide):
    """Citizen Portal."""
    _feature_list_with_screenshot(slide, SLIDE5["title"], SLIDE5["features"],
                                  SLIDE5["screenshot_label"], "citizen_dashboard")


def build_slide6(slide):
    """Official Portal."""
    _feature_list_with_screenshot(slide, SLIDE6["title"], SLIDE6["features"],
                                  SLIDE6["screenshot_label"], "official_dashboard")


def build_slide7(slide):
    """Admin Portal."""
    _feature_list_with_screenshot(slide, SLIDE7["title"], SLIDE7["features"],
                                  SLIDE7["screenshot_label"], "admin_dashboard")


def build_slide8(slide):
    """Continuous Governance Loop."""
    clear_slide_text(slide)
    _title_on_template(slide, SLIDE8["title"])

    # Governance loop diagram — reduced height to Inches(1.8)
    gl = diagram_path("governance_loop")
    if gl:
        slide.shapes.add_picture(gl, Inches(0.3), Inches(1.4), Inches(9.4), Inches(1.8))

    # Step descriptions below diagram to fill remaining space
    desc_colors = [GREEN_PRIMARY, BLUE_LIGHT, ORANGE_ACCENT,
                   RGBColor(0x7B, 0x1F, 0xA2), RED_URGENT]
    y = Inches(3.3)
    for i, (step_name, desc) in enumerate(SLIDE8.get("descriptions", [])):
        color = desc_colors[i % len(desc_colors)]
        # Colored badge
        add_flow_box(slide, Inches(0.3), y, Inches(1.3), Inches(0.28),
                     step_name, color, WHITE, Pt(8))
        # Description
        tf_d = add_textbox(slide, Inches(1.7), y, Inches(8.0), Inches(0.28))
        run = tf_d.paragraphs[0].add_run()
        run.text = desc
        set_font(run, Pt(7.5), color=DARK_TEXT)
        y += Inches(0.33)

    # Caption
    tf_c = add_textbox(slide, Inches(0.3), Inches(5.0), Inches(9.4), Inches(0.3))
    tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf_c.paragraphs[0].add_run()
    run.text = SLIDE8["caption"]
    set_font(run, Pt(8), color=GRAY_MED, italic=True)


def build_slide9(slide):
    """Why AI? Before vs After."""
    clear_slide_text(slide)
    d = SLIDE9
    _title_on_template(slide, d["title"])

    table = add_table(slide, len(d["comparisons"]) + 1, 3,
                      Inches(0.3), Inches(1.4), Inches(9.4), Inches(3.8))
    set_cell_text(table, 0, 0, "User", SIZE_TABLE_HEADER, True, WHITE)
    set_cell_text(table, 0, 1, "Without AI", SIZE_TABLE_HEADER, True, WHITE)
    set_cell_text(table, 0, 2, "With eGramSabha AI", SIZE_TABLE_HEADER, True, WHITE)
    for i, (user, without, with_ai) in enumerate(d["comparisons"]):
        set_cell_text(table, i + 1, 0, user, Pt(7.5), True)
        set_cell_text(table, i + 1, 1, without, Pt(7.5), color=RED_URGENT)
        set_cell_text(table, i + 1, 2, with_ai, Pt(7.5), color=GREEN_PRIMARY)
    style_table(table)
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(3.7)
    table.columns[2].width = Inches(4.2)


def build_slide10(slide):
    """Voice Issue Reporting Pipeline."""
    clear_slide_text(slide)
    d = SLIDE10
    _title_on_template(slide, d["title"])

    # Diagram — reduced height
    vp = diagram_path("voice_pipeline")
    if vp:
        slide.shapes.add_picture(vp, Inches(0.2), Inches(1.4), Inches(9.6), Inches(1.2))
        y_start = Inches(2.7)
    else:
        y_start = Inches(1.5)

    # Steps as compact numbered list
    tf = add_textbox(slide, Inches(0.3), y_start, Inches(9.4), Inches(5.0) - y_start)
    for i, (step, detail) in enumerate(d["steps"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = f"{i + 1}. {step}: "
        set_font(run, Pt(9), True, GREEN_PRIMARY)
        run2 = p.add_run()
        run2.text = detail
        set_font(run2, Pt(8), color=DARK_TEXT)


def build_slide11(slide):
    """AI Issue Summarisation & Agenda."""
    clear_slide_text(slide)
    d = SLIDE11
    _title_on_template(slide, d["title"])

    # Left column: Summarization with colored header
    add_flow_box(slide, Inches(0.3), Inches(1.4), Inches(4.5), Inches(0.35),
                 "AI Summarisation (Hourly Cron)", GREEN_PRIMARY, WHITE, Pt(10))

    tf_l = add_textbox(slide, Inches(0.3), Inches(1.8), Inches(4.5), Inches(1.8))
    for i, step in enumerate(d["summarization_steps"]):
        p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
        p.space_before = Pt(2)
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = f"{i + 1}. {step}"
        set_font(run, Pt(7.5), color=DARK_TEXT)

    # Arrow between columns (centered in the gap between left @4.8" and right @5.1")
    add_arrow(slide, Inches(4.82), Inches(2.3), Inches(0.25))

    # Right column: Agenda with colored header
    add_flow_box(slide, Inches(5.1), Inches(1.4), Inches(4.5), Inches(0.35),
                 "Agenda Generation (Official-Driven)", BLUE_LIGHT, WHITE, Pt(10))

    tf_r = add_textbox(slide, Inches(5.1), Inches(1.8), Inches(4.5), Inches(2.2))
    for i, step in enumerate(d["agenda_steps"]):
        p = tf_r.paragraphs[0] if i == 0 else tf_r.add_paragraph()
        p.space_before = Pt(2)
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = f"{i + 1}. {step}"
        set_font(run, Pt(7.5), color=DARK_TEXT)

    # Key flow text between the two sections
    if "key_flow" in d:
        tf_flow = add_textbox(slide, Inches(0.3), Inches(3.85), Inches(9.4), Inches(0.3))
        tf_flow.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf_flow.paragraphs[0].add_run()
        run.text = d["key_flow"]
        set_font(run, Pt(7), color=GRAY_MED, italic=True, name="Consolas")

    # Key point callout at bottom
    add_callout_box(slide, Inches(0.3), Inches(4.2), Inches(9.4), Inches(0.45),
                    d["key_point"], GREEN_PRIMARY, Pt(10))


def build_slide12(slide):
    """Face Authentication."""
    clear_slide_text(slide)
    d = SLIDE12
    _title_on_template(slide, d["title"])

    # Face auth diagram — reduced height
    fa = diagram_path("face_auth")
    if fa:
        slide.shapes.add_picture(fa, Inches(0.2), Inches(1.4), Inches(9.6), Inches(1.0))
        y_start = Inches(2.5)
    else:
        y_start = Inches(1.4)

    # Left: Registration
    add_flow_box(slide, Inches(0.3), y_start, Inches(4.5), Inches(0.3),
                 "Registration Flow", GREEN_PRIMARY, WHITE, Pt(10))
    tf_l = add_textbox(slide, Inches(0.3), y_start + Inches(0.35), Inches(4.5), Inches(1.8))
    for i, step in enumerate(d["registration"]):
        p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = f"{i + 1}. {step}"
        set_font(run, Pt(8), color=DARK_TEXT)

    # Right: Login
    add_flow_box(slide, Inches(5.1), y_start, Inches(4.5), Inches(0.3),
                 "Login Flow", BLUE_LIGHT, WHITE, Pt(10))
    tf_r = add_textbox(slide, Inches(5.1), y_start + Inches(0.35), Inches(4.5), Inches(2.0))
    for i, step in enumerate(d["login"]):
        p = tf_r.paragraphs[0] if i == 0 else tf_r.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = f"{i + 1}. {step}"
        set_font(run, Pt(8), color=DARK_TEXT)

    # Production note
    tf_p = add_textbox(slide, Inches(0.3), Inches(4.9), Inches(9.4), Inches(0.3))
    run = tf_p.paragraphs[0].add_run()
    run.text = f"Production: {d['production']}"
    set_font(run, Pt(7), color=GRAY_MED, italic=True)


def build_slide13(slide):
    """Meeting Day Flow."""
    clear_slide_text(slide)
    d = SLIDE13
    _title_on_template(slide, d["title"])

    # Meeting day diagram
    md = diagram_path("meeting_day")
    if md:
        slide.shapes.add_picture(md, Inches(0.2), Inches(1.4), Inches(9.6), Inches(1.2))
        y_start = Inches(2.7)
    else:
        y_start = Inches(1.4)

    # Steps as two columns
    steps = d["steps"]
    mid = 3
    col_w = Inches(4.5)

    for col, (start, end, x) in enumerate([(0, mid, Inches(0.3)), (mid, len(steps), Inches(5.1))]):
        tf = add_textbox(slide, x, y_start, col_w, Inches(2.0))
        for i, (step_name, desc) in enumerate(steps[start:end]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(3)
            p.space_after = Pt(1)
            run = p.add_run()
            run.text = f"{start + i + 1}. {step_name}: "
            set_font(run, Pt(9), True, GREEN_PRIMARY)
            run2 = p.add_run()
            run2.text = desc
            set_font(run2, Pt(8), color=DARK_TEXT)

    # MOM planned callout
    add_callout_box(slide, Inches(0.3), Inches(4.3), Inches(9.4), Inches(0.4),
                    d["planned"], ORANGE_ACCENT, Pt(8))


# ============================================================
# EXTRA SLIDE BUILDERS (14-30)
# ============================================================

def build_slide14(prs):
    """TTS & Multilingual Support."""
    slide = add_extra_slide(prs, SLIDE14["title"])
    d = SLIDE14

    # Left: TTS Flow with colored step boxes
    add_flow_box(slide, Inches(0.3), Inches(1.5), Inches(5.0), Inches(0.3),
                 "TTS Caching Flow (Amazon Polly)", GREEN_PRIMARY, WHITE, Pt(10))

    step_colors = [BLUE_LIGHT, GREEN_PRIMARY, ORANGE_ACCENT,
                   GREEN_PRIMARY, RGBColor(0x7B, 0x1F, 0xA2), BLUE_LIGHT]
    y = Inches(1.9)
    for i, step in enumerate(d["tts_flow"]):
        # Small colored indicator
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.35), y + Inches(0.02),
                                       Inches(0.18), Inches(0.18))
        shape.fill.solid()
        shape.fill.fore_color.rgb = step_colors[i % len(step_colors)]
        shape.line.fill.background()

        tf_s = add_textbox(slide, Inches(0.6), y, Inches(4.7), Inches(0.25))
        run = tf_s.paragraphs[0].add_run()
        run.text = f"{i + 1}. {step}"
        set_font(run, Pt(8), color=DARK_TEXT)
        y += Inches(0.28)

    # Right: Languages as grid
    add_flow_box(slide, Inches(5.5), Inches(1.5), Inches(4.2), Inches(0.3),
                 "Supported Languages (10+)", BLUE_LIGHT, WHITE, Pt(10))

    for i, lang in enumerate(d["languages"]):
        row, col = i // 2, i % 2
        x = Inches(5.5) + col * Inches(2.1)
        y = Inches(1.95) + row * Inches(0.32)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x, y, Inches(1.9), Inches(0.28))
        shape.fill.solid()
        shape.fill.fore_color.rgb = TABLE_ALT_BG
        shape.line.fill.background()
        tf_lang = shape.text_frame
        tf_lang.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf_lang.paragraphs[0].add_run()
        run.text = lang
        set_font(run, Pt(9), True, GREEN_PRIMARY)

    # Bottom notes
    tf_ui = add_textbox(slide, Inches(0.3), Inches(4.5), Inches(4.5), Inches(0.5))
    run = tf_ui.paragraphs[0].add_run()
    run.text = f"UI: {d['ui_i18n']}"
    set_font(run, Pt(8), color=DARK_TEXT)

    tf_ai = add_textbox(slide, Inches(5.0), Inches(4.5), Inches(4.7), Inches(0.5))
    run = tf_ai.paragraphs[0].add_run()
    run.text = f"AI Content: {d['ai_translation']}"
    set_font(run, Pt(8), color=DARK_TEXT)

    return slide


def build_slide15(prs):
    """System Architecture - Diagram."""
    slide = add_extra_slide(prs, "")  # title manually placed higher
    d = SLIDE15
    # Title at 0.70" for more vertical space
    tf = add_textbox(slide, Inches(0.3), Inches(0.70), CONTENT_WIDTH, Inches(0.5))
    run = tf.paragraphs[0].add_run()
    run.text = d["title"]
    set_font(run, SIZE_SLIDE_TITLE, bold=True, color=DARK_TEXT)

    arch = diagram_path("architecture")
    if arch:
        slide.shapes.add_picture(arch, Inches(0.3), Inches(1.12), Inches(9.4), Inches(3.8))
    else:
        add_placeholder(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(3.5),
                        "[Architecture Diagram]\n6-Layer System Architecture")

    # Bottom info bar
    info_items = [
        (d["docker"], GREEN_PRIMARY),
        (d["services_active"], BLUE_LIGHT),
        (d["services_ready"], ORANGE_ACCENT),
    ]
    for i, (text, color) in enumerate(info_items):
        x = Inches(0.3) + i * Inches(3.2)
        add_flow_box(slide, x, Inches(5.1), Inches(3.0), Inches(0.35),
                     text, color, WHITE, Pt(8))

    return slide


def build_slide16(prs):
    """Architecture Layers - Table."""
    slide = add_extra_slide(prs, SLIDE16["title"])
    d = SLIDE16

    table = add_table(slide, len(d["layers"]) + 1, 2,
                      Inches(0.3), Inches(1.5), Inches(9.4), Inches(3.5))
    set_cell_text(table, 0, 0, "Layer", SIZE_TABLE_HEADER, True, WHITE)
    set_cell_text(table, 0, 1, "Components & Details", SIZE_TABLE_HEADER, True, WHITE)

    layer_icons = ["🌐", "🔗", "🛡️", "⚙️", "🤖", "💾"]
    for i, (layer, components) in enumerate(d["layers"]):
        set_cell_text(table, i + 1, 0, layer, Pt(9), True)
        set_cell_text(table, i + 1, 1, components, Pt(8))

    style_table(table)
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(6.9)

    return slide


def build_slide17(prs):
    """Security Architecture."""
    slide = add_extra_slide(prs, "")  # title manually placed higher
    d = SLIDE17
    # Title at 0.70" for more vertical space
    tf = add_textbox(slide, Inches(0.3), Inches(0.70), CONTENT_WIDTH, Inches(0.5))
    run = tf.paragraphs[0].add_run()
    run.text = d["title"]
    set_font(run, SIZE_SLIDE_TITLE, bold=True, color=DARK_TEXT)

    layer_colors = [RED_URGENT, ORANGE_ACCENT, BLUE_LIGHT,
                    GREEN_PRIMARY, RGBColor(0x7B, 0x1F, 0xA2), BLUE_AWS]
    y = Inches(1.26)
    for i, ((layer_name, desc), color) in enumerate(zip(d["layers"], layer_colors)):
        # Layer badge
        add_flow_box(slide, Inches(0.3), y, Inches(2.5), Inches(0.5),
                     layer_name, color, WHITE, Pt(9))
        # Description
        tf_d = add_textbox(slide, Inches(3.0), y + Inches(0.07), Inches(6.7), Inches(0.45))
        run = tf_d.paragraphs[0].add_run()
        run.text = desc
        set_font(run, Pt(8), color=DARK_TEXT)
        y += Inches(0.55)

    # RBAC roles
    add_section_header(slide, "RBAC Roles", Inches(0.3), Inches(4.7), Inches(2.0))
    role_x = Inches(0.3)
    for i, role in enumerate(d["roles"]):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       role_x, Inches(5.0), Inches(1.2), Inches(0.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = TABLE_ALT_BG
        shape.line.fill.background()
        tf_r = shape.text_frame
        tf_r.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf_r.paragraphs[0].add_run()
        run.text = role
        set_font(run, Pt(7), True, GREEN_PRIMARY)
        role_x += Inches(1.3)

    return slide


def build_slide18(prs):
    """Traffic Flow."""
    slide = add_extra_slide(prs, SLIDE18["title"])
    d = SLIDE18

    flow_colors = [GREEN_PRIMARY, BLUE_LIGHT, ORANGE_ACCENT]
    y = Inches(1.55)
    for i, ((flow_name, path, desc), color) in enumerate(zip(d["flows"], flow_colors)):
        # Flow name badge
        add_flow_box(slide, Inches(0.3), y, Inches(1.5), Inches(0.7),
                     flow_name, color, WHITE, Pt(10))
        # Path
        tf_p = add_textbox(slide, Inches(2.0), y, Inches(7.7), Inches(0.35))
        run = tf_p.paragraphs[0].add_run()
        run.text = path
        set_font(run, Pt(8), True, DARK_TEXT)
        # Description
        tf_d = add_textbox(slide, Inches(2.0), y + Inches(0.35), Inches(7.7), Inches(0.35))
        run = tf_d.paragraphs[0].add_run()
        run.text = desc
        set_font(run, Pt(7.5), color=GRAY_MED)
        y += Inches(0.85)

    # Docker volumes
    add_section_header(slide, "Docker Volumes", Inches(0.3), Inches(4.15), Inches(3.0))
    for i, (vol_name, vol_desc) in enumerate(d["volumes"]):
        x = Inches(0.3) + i * Inches(3.2)
        # Volume box with name and description
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x, Inches(4.5), Inches(3.0), Inches(0.55))
        shape.fill.solid()
        shape.fill.fore_color.rgb = TABLE_ALT_BG
        shape.line.color.rgb = GREEN_LIGHT
        shape.line.width = Pt(1)
        tf_v = shape.text_frame
        tf_v.word_wrap = True
        tf_v.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf_v.paragraphs[0].add_run()
        run.text = vol_name
        set_font(run, Pt(8), True, GREEN_PRIMARY)
        p2 = tf_v.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = vol_desc
        set_font(run2, Pt(7), color=GRAY_MED)

    return slide


def build_slide19(prs):
    """AWS Services (15)."""
    slide = add_extra_slide(prs, "")  # title manually placed higher
    d = SLIDE19
    # Title at 0.66" for more vertical space (dense slide)
    tf = add_textbox(slide, CONTENT_LEFT, Inches(0.66), CONTENT_WIDTH, Inches(0.5))
    run = tf.paragraphs[0].add_run()
    run.text = d["title"]
    set_font(run, SIZE_SLIDE_TITLE, bold=True, color=DARK_TEXT)

    # Active services table
    n_active = len(d["active"])
    table = add_table(slide, n_active + 1, 3,
                      Inches(0.3), Inches(1.21), Inches(9.4), Inches(3.02))
    set_cell_text(table, 0, 0, "Service", SIZE_TABLE_HEADER, True, WHITE)
    set_cell_text(table, 0, 1, "Configuration", SIZE_TABLE_HEADER, True, WHITE)
    set_cell_text(table, 0, 2, "Usage", SIZE_TABLE_HEADER, True, WHITE)
    for i, (svc, config, usage) in enumerate(d["active"]):
        set_cell_text(table, i + 1, 0, svc, Pt(7), True)
        set_cell_text(table, i + 1, 1, config, Pt(7))
        set_cell_text(table, i + 1, 2, usage, Pt(7))
    style_table(table)
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(3.0)
    table.columns[2].width = Inches(4.6)

    # Code-ready section
    table_bottom = Inches(1.21) + Inches(3.02)  # table top + table height
    y_cr = table_bottom + Inches(0.15)
    add_flow_box(slide, Inches(0.3), y_cr, Inches(2.5), Inches(0.25),
                 "Code-Ready (4 Services)", YELLOW_READY, WHITE, Pt(8))

    y_cr += Inches(0.3)
    cr_colors = [ORANGE_ACCENT, BLUE_LIGHT, GREEN_PRIMARY, RGBColor(0x7B, 0x1F, 0xA2)]
    for i, ((svc, desc), color) in enumerate(zip(d["code_ready"], cr_colors)):
        col = i % 2
        row = i // 2
        x = Inches(0.3) + col * Inches(4.8)
        y = y_cr + row * Inches(0.4)

        # Colored left accent bar + text
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        x, y, Inches(0.08), Inches(0.35))
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()

        tf_s = add_textbox(slide, x + Inches(0.12), y, Inches(4.5), Inches(0.35))
        run = tf_s.paragraphs[0].add_run()
        run.text = svc
        set_font(run, Pt(7), True, color)
        p2 = tf_s.add_paragraph()
        p2.space_before = Pt(0)
        run2 = p2.add_run()
        run2.text = desc
        set_font(run2, Pt(6.5), color=DARK_TEXT)

    return slide


def build_slide20(prs):
    """Provider Abstraction."""
    slide = add_extra_slide(prs, SLIDE20["title"])
    d = SLIDE20

    table = add_table(slide, len(d["switches"]) + 1, 3,
                      Inches(0.3), Inches(1.5), Inches(9.4), Inches(2.5))
    set_cell_text(table, 0, 0, "Environment Variable", SIZE_TABLE_HEADER, True, WHITE)
    set_cell_text(table, 0, 1, "Options", SIZE_TABLE_HEADER, True, WHITE)
    set_cell_text(table, 0, 2, "Purpose", SIZE_TABLE_HEADER, True, WHITE)
    for i, (var, options, purpose) in enumerate(d["switches"]):
        set_cell_text(table, i + 1, 0, var, Pt(8), True)
        set_cell_text(table, i + 1, 1, options, Pt(8))
        set_cell_text(table, i + 1, 2, purpose, Pt(8))
    style_table(table)
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(3.2)
    table.columns[2].width = Inches(3.2)

    add_callout_box(slide, Inches(0.3), Inches(4.3), Inches(9.4), Inches(0.45),
                    d["benefit"], GREEN_PRIMARY, Pt(10))

    return slide


def build_slide21(prs):
    """Technologies Utilized."""
    slide = add_extra_slide(prs, SLIDE21["title"])
    d = SLIDE21

    # 4 quadrant cards
    positions = [
        (Inches(0.3), Inches(1.5), Inches(4.5), Inches(1.65)),   # Frontend (top-left)
        (Inches(5.0), Inches(1.5), Inches(4.7), Inches(1.65)),   # Backend (top-right)
        (Inches(0.3), Inches(3.45), Inches(4.5), Inches(1.65)),   # AI Backend (bottom-left)
        (Inches(5.0), Inches(3.45), Inches(4.7), Inches(1.65)),   # Database (bottom-right)
    ]

    for idx, (section_name, color_hex, items) in enumerate(d["sections"]):
        x, y, w, h = positions[idx]
        color = RGBColor.from_string(color_hex[1:]) if color_hex.startswith("#") else GREEN_PRIMARY

        # Section header bar — slightly larger with underline accent
        hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.33))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        tf_h = hdr.text_frame
        tf_h.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf_h.paragraphs[0].add_run()
        run.text = section_name
        set_font(run, Pt(11), True, WHITE)

        # Items with colored left accent bars
        item_y = y + Inches(0.4)
        for i, (tech, purpose) in enumerate(items):
            # Colored left accent bar
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            x + Inches(0.05), item_y,
                                            Inches(0.06), Inches(0.24))
            accent.fill.solid()
            accent.fill.fore_color.rgb = color
            accent.line.fill.background()

            # Tech name (extract version if present for badge)
            tf_tech = add_textbox(slide, x + Inches(0.15), item_y - Inches(0.02),
                                  w - Inches(0.2), Inches(0.26))
            run = tf_tech.paragraphs[0].add_run()
            run.text = f"{tech}: "
            set_font(run, Pt(7.5), True, color)
            run2 = tf_tech.paragraphs[0].add_run()
            run2.text = purpose
            set_font(run2, Pt(7), color=DARK_TEXT)

            item_y += Inches(0.27)

    return slide


def build_slide22(prs):
    """Cost Breakdown."""
    slide = add_extra_slide(prs, SLIDE22["title"])
    d = SLIDE22

    # Full-width cost table
    table = add_table(slide, len(d["costs"]) + 2, 3,
                      Inches(0.3), Inches(1.5), Inches(9.4), Inches(3.5))
    set_cell_text(table, 0, 0, "Service", SIZE_TABLE_HEADER, True, WHITE)
    set_cell_text(table, 0, 1, "Monthly Cost", SIZE_TABLE_HEADER, True, WHITE)
    set_cell_text(table, 0, 2, "Notes", SIZE_TABLE_HEADER, True, WHITE)
    for i, (svc, cost, notes) in enumerate(d["costs"]):
        set_cell_text(table, i + 1, 0, svc, Pt(9))
        set_cell_text(table, i + 1, 1, cost, Pt(9), True)
        set_cell_text(table, i + 1, 2, notes, Pt(8))
    last = len(d["costs"]) + 1
    set_cell_text(table, last, 0, "TOTAL", Pt(10), True, GREEN_PRIMARY)
    set_cell_text(table, last, 1, d["total"], Pt(10), True, GREEN_PRIMARY)
    set_cell_text(table, last, 2, "Hackathon scale (1-5 panchayats)", Pt(8))
    style_table(table)
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(4.9)

    return slide


def build_slide23(prs):
    """Cost Optimizations."""
    slide = add_extra_slide(prs, SLIDE23["title"])
    d = SLIDE23

    y = Inches(1.5)
    for i, (name, desc, color_hex) in enumerate(d["optimizations"]):
        color = RGBColor.from_string(color_hex[1:])

        # Left: optimization name badge
        add_flow_box(slide, Inches(0.3), y, Inches(2.8), Inches(0.55),
                     name, color, WHITE, Pt(9))

        # Right: description with accent bar
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(3.3), y, Inches(0.06), Inches(0.55))
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()

        tf_d = add_textbox(slide, Inches(3.5), y + Inches(0.08), Inches(6.2), Inches(0.45))
        run = tf_d.paragraphs[0].add_run()
        run.text = desc
        set_font(run, Pt(9), color=DARK_TEXT)

        y += Inches(0.65)

    return slide


def build_slide24(prs):
    """Scale Projections."""
    slide = add_extra_slide(prs, SLIDE24["title"])
    d = SLIDE24

    table = add_table(slide, len(d["scale"]) + 1, 5,
                      Inches(0.3), Inches(1.5), Inches(9.4), Inches(1.6))
    for c, h in enumerate(["Phase", "Scale", "Changes Required", "Monthly Cost", "Per Citizen/Yr"]):
        set_cell_text(table, 0, c, h, SIZE_TABLE_HEADER, True, WHITE)
    for i, (phase, scale, changes, cost, per_citizen) in enumerate(d["scale"]):
        set_cell_text(table, i + 1, 0, phase, Pt(8), True)
        set_cell_text(table, i + 1, 1, scale, Pt(7.5))
        set_cell_text(table, i + 1, 2, changes, Pt(7))
        set_cell_text(table, i + 1, 3, cost, Pt(8), True)
        pc_color = GREEN_PRIMARY if "$" in per_citizen else DARK_TEXT
        set_cell_text(table, i + 1, 4, per_citizen, Pt(8), True, pc_color)
    style_table(table)
    table.columns[0].width = Inches(1.1)
    table.columns[1].width = Inches(1.4)
    table.columns[2].width = Inches(3.7)
    table.columns[3].width = Inches(1.2)
    table.columns[4].width = Inches(2.0)

    # Comparison callout
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.3), Inches(3.35), Inches(9.4), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN_PRIMARY
    shape.line.fill.background()
    tf_c = shape.text_frame
    tf_c.word_wrap = True
    tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf_c.paragraphs[0].add_run()
    run.text = d["reduction"]
    set_font(run, Pt(16), True, WHITE)
    p2 = tf_c.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = d["comparison"]
    set_font(run2, Pt(9), color=WHITE)

    # National note
    tf_n = add_textbox(slide, Inches(0.3), Inches(4.2), Inches(9.4), Inches(0.4))
    tf_n.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf_n.paragraphs[0].add_run()
    run.text = d["national_note"]
    set_font(run, Pt(9), color=GRAY_MED, italic=True)

    return slide


def build_slide25(prs):
    """Impact Potential."""
    slide = add_extra_slide(prs, SLIDE25["title"])
    d = SLIDE25

    table = add_table(slide, len(d["metrics"]) + 1, 4,
                      Inches(0.3), Inches(1.5), Inches(9.4), Inches(2.0))
    set_cell_text(table, 0, 0, "Metric", SIZE_TABLE_HEADER, True, WHITE)
    set_cell_text(table, 0, 1, "National Average", SIZE_TABLE_HEADER, True, WHITE)
    set_cell_text(table, 0, 2, "With Digital Gram Sabha", SIZE_TABLE_HEADER, True, WHITE)
    set_cell_text(table, 0, 3, "Improvement", SIZE_TABLE_HEADER, True, WHITE)
    for i, (metric, avg, digital, improvement) in enumerate(d["metrics"]):
        set_cell_text(table, i + 1, 0, metric, Pt(8))
        set_cell_text(table, i + 1, 1, avg, Pt(8), color=RED_URGENT)
        set_cell_text(table, i + 1, 2, digital, Pt(8), True, GREEN_PRIMARY)
        set_cell_text(table, i + 1, 3, improvement, Pt(9), True, GREEN_PRIMARY)
    style_table(table)
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(2.3)
    table.columns[2].width = Inches(3.0)
    table.columns[3].width = Inches(2.3)

    # Approach
    tf_a = add_textbox(slide, Inches(0.3), Inches(3.7), Inches(9.4), Inches(0.7))
    run = tf_a.paragraphs[0].add_run()
    run.text = "Registration Approach: "
    set_font(run, Pt(9), True, GREEN_PRIMARY)
    run2 = tf_a.paragraphs[0].add_run()
    run2.text = d["approach"]
    set_font(run2, Pt(8), color=DARK_TEXT)

    return slide


def build_slide26(prs):
    """Impact Outcomes."""
    slide = add_extra_slide(prs, SLIDE26["title"])
    d = SLIDE26

    outcome_colors = [GREEN_PRIMARY, BLUE_LIGHT, ORANGE_ACCENT, RGBColor(0x7B, 0x1F, 0xA2),
                      RED_URGENT, GREEN_PRIMARY, BLUE_LIGHT, ORANGE_ACCENT]

    # Two columns of outcomes with colored accent bars
    mid = 4
    for col, (start, end, x) in enumerate([(0, mid, Inches(0.3)), (mid, len(d["outcomes"]), Inches(5.0))]):
        y = Inches(1.5)
        for i, (name, desc) in enumerate(d["outcomes"][start:end]):
            idx = start + i
            color = outcome_colors[idx % len(outcome_colors)]

            # Colored accent bar
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            x, y, Inches(0.08), Inches(0.55))
            accent.fill.solid()
            accent.fill.fore_color.rgb = color
            accent.line.fill.background()

            # Name + description
            tf = add_textbox(slide, x + Inches(0.15), y, Inches(4.5), Inches(0.55))
            run = tf.paragraphs[0].add_run()
            run.text = name
            set_font(run, Pt(9), True, color)
            p2 = tf.add_paragraph()
            p2.space_before = Pt(1)
            run2 = p2.add_run()
            run2.text = desc
            set_font(run2, Pt(7.5), color=DARK_TEXT)

            y += Inches(0.6)

    # Takeaway callout
    add_callout_box(slide, Inches(0.3), Inches(4.8), Inches(9.4), Inches(0.4),
                    d["takeaway"], GREEN_PRIMARY, Pt(8))

    return slide


def build_slide27(prs):
    """National Impact Projection."""
    slide = add_extra_slide(prs, SLIDE27["title"])
    d = SLIDE27

    table = add_table(slide, len(d["impact"]) + 1, 3,
                      Inches(0.3), Inches(1.5), Inches(9.4), Inches(1.8))
    set_cell_text(table, 0, 0, "Impact Area", SIZE_TABLE_HEADER, True, WHITE)
    set_cell_text(table, 0, 1, "Mechanism", SIZE_TABLE_HEADER, True, WHITE)
    set_cell_text(table, 0, 2, "Projected Outcome", SIZE_TABLE_HEADER, True, WHITE)
    for i, (area, mechanism, outcome) in enumerate(d["impact"]):
        set_cell_text(table, i + 1, 0, area, Pt(8), True)
        set_cell_text(table, i + 1, 1, mechanism, Pt(7.5))
        set_cell_text(table, i + 1, 2, outcome, Pt(8), True, GREEN_PRIMARY)
    style_table(table)
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(4.0)
    table.columns[2].width = Inches(3.6)

    # National scale
    tf_n = add_textbox(slide, Inches(0.3), Inches(3.5), Inches(9.4), Inches(0.6))
    run = tf_n.paragraphs[0].add_run()
    run.text = d["national_scale"]
    set_font(run, Pt(8), color=DARK_TEXT, italic=True)

    # Multiplier effects as horizontal cards
    add_section_header(slide, "Multiplier Effects", Inches(0.3), Inches(4.2), Inches(3.0))
    for i, effect in enumerate(d["multipliers"]):
        col = i % 2
        row = i // 2
        x = Inches(0.3) + col * Inches(4.8)
        y = Inches(4.5) + row * Inches(0.35)

        # Dot indicator
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.05),
                                     Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GREEN_PRIMARY
        dot.line.fill.background()

        tf_e = add_textbox(slide, x + Inches(0.2), y, Inches(4.4), Inches(0.3))
        run = tf_e.paragraphs[0].add_run()
        run.text = effect
        set_font(run, Pt(8), color=DARK_TEXT)

    return slide


def _build_gallery_slide(prs, data):
    """Build a single-screenshot gallery slide with context panel on the left."""
    slide = add_extra_slide(prs, data["title"])
    portal_color = RGBColor.from_string(data["portal_color"][1:])

    # --- Left panel: context (3.4" wide) ---
    panel_x = CONTENT_LEFT
    panel_w = Inches(3.4)
    body_top = Inches(1.5)

    # Portal badge
    add_flow_box(slide, panel_x, body_top, Inches(2.2), Inches(0.32),
                 data["portal"], portal_color, WHITE, Pt(9))

    # Page name
    tf_name = add_textbox(slide, panel_x, body_top + Inches(0.45), panel_w, Inches(0.35))
    run = tf_name.paragraphs[0].add_run()
    run.text = data["page_name"]
    set_font(run, Pt(14), True, DARK_TEXT)

    # Thin accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  panel_x, body_top + Inches(0.82), panel_w, Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = portal_color
    line.line.fill.background()

    # Description
    tf_desc = add_textbox(slide, panel_x, body_top + Inches(0.92), panel_w, Inches(1.2))
    run = tf_desc.paragraphs[0].add_run()
    run.text = data["description"]
    set_font(run, Pt(8), color=DARK_TEXT)

    # Key highlights header
    tf_hdr = add_textbox(slide, panel_x, body_top + Inches(2.15), panel_w, Inches(0.25))
    run = tf_hdr.paragraphs[0].add_run()
    run.text = "Key Highlights"
    set_font(run, Pt(9), True, portal_color)

    # Highlight bullets with colored dots
    y_hl = body_top + Inches(2.42)
    for hl in data["highlights"]:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     panel_x + Inches(0.02), y_hl + Inches(0.04),
                                     Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = portal_color
        dot.line.fill.background()

        tf_hl = add_textbox(slide, panel_x + Inches(0.18), y_hl, panel_w - Inches(0.18), Inches(0.35))
        run = tf_hl.paragraphs[0].add_run()
        run.text = hl
        set_font(run, Pt(7.5), color=DARK_TEXT)
        y_hl += Inches(0.32)

    # --- Right panel: screenshot area ---
    area_x = Inches(3.9)
    area_y = body_top
    area_w = Inches(5.8)
    area_h = Inches(3.7)

    sp = screenshot_path(data["screenshot"])
    if sp:
        try:
            w, h = fit_image(sp, area_w, area_h)
            # Center image within the area
            ix = area_x + (area_w - w) // 2
            iy = area_y + (area_h - h) // 2

            # Sharp rectangle border aligned exactly to image
            border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            ix - Inches(0.04), iy - Inches(0.04),
                                            w + Inches(0.08), h + Inches(0.08))
            border.fill.solid()
            border.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
            border.line.color.rgb = RGBColor(0xBD, 0xBD, 0xBD)
            border.line.width = Pt(1)

            slide.shapes.add_picture(sp, ix, iy, w, h)
            return slide
        except Exception:
            pass

    add_placeholder(slide, area_x, area_y, area_w, area_h,
                    f"[Screenshot]\n{data['page_name']}")

    return slide


def build_gallery_slides(prs):
    """Screenshot Gallery."""
    slides = []
    for data in GALLERY_SLIDES:
        slides.append(_build_gallery_slide(prs, data))
    return slides


def build_slide30(prs):
    """Future Roadmap."""
    slide = add_extra_slide(prs, SLIDE30["title"])
    d = SLIDE30

    col_width = Inches(2.2)
    gap = Inches(0.12)

    for i, (phase_name, color_hex, items) in enumerate(d["phases"]):
        color = RGBColor.from_string(color_hex[1:])
        x = Inches(0.3) + i * (col_width + gap)
        y = Inches(1.5)

        # Phase header
        add_flow_box(slide, x, y, col_width, Inches(0.5),
                     phase_name, color, WHITE, Pt(8))

        # Connecting arrow (except last)
        if i < 3:
            arrow_x = x + col_width + Inches(0.01)
            add_arrow(slide, arrow_x, y + Inches(0.15), gap - Inches(0.02), Inches(0.2))

        # Items with accent dots
        y_item = y + Inches(0.6)
        for j, item in enumerate(items):
            # Small dot
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         x + Inches(0.05), y_item + Inches(0.04),
                                         Inches(0.1), Inches(0.1))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()

            tf_item = add_textbox(slide, x + Inches(0.2), y_item,
                                  col_width - Inches(0.25), Inches(0.3))
            run = tf_item.paragraphs[0].add_run()
            run.text = item
            set_font(run, Pt(7), color=DARK_TEXT)
            y_item += Inches(0.3)

    return slide


def build_closing_slide(prs):
    """Our Thank You / Closing (extra slide)."""
    slide = add_extra_slide(prs, "")
    d = SLIDE31

    # Logo centered (below template header area)
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(4.0), Inches(0.6), Inches(1.8), Inches(1.8))

    # eGramSabha name
    tf = add_textbox(slide, Inches(0.5), Inches(2.5), Inches(9.0), Inches(0.5))
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = d["closing"]
    set_font(run, Pt(20), True, GREEN_PRIMARY)

    # Vision
    tf2 = add_textbox(slide, Inches(0.5), Inches(3.1), Inches(9.0), Inches(0.5))
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf2.paragraphs[0].add_run()
    run.text = d["vision"]
    set_font(run, Pt(16), True, DARK_TEXT, italic=True)

    # Subtitle
    tf3 = add_textbox(slide, Inches(0.5), Inches(3.6), Inches(9.0), Inches(0.4))
    tf3.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf3.paragraphs[0].add_run()
    run.text = d["subtitle"]
    set_font(run, Pt(12), color=GRAY_MED)

    # Links row
    links = SLIDE13_LINKS
    link_items = []
    if links["github_repo"]:
        link_items.append(("GitHub", links["github_repo"]))
    if links["demo_video"]:
        link_items.append(("Demo Video", links["demo_video"]))
    if links["live_prototype"]:
        link_items.append(("Live Demo", links["live_prototype"]))

    tf_links = add_textbox(slide, Inches(0.5), Inches(3.85), Inches(9.0), Inches(0.9))
    for idx, (label, url) in enumerate(link_items):
        p = tf_links.paragraphs[0] if idx == 0 else tf_links.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        run_label = p.add_run()
        run_label.text = f"{label}: "
        set_font(run_label, Pt(9), bold=True, color=DARK_TEXT)
        run_url = p.add_run()
        run_url.text = url
        set_font(run_url, Pt(9), color=BLUE_LIGHT)
        run_url.hyperlink.address = url

    # Credentials
    cred_items = [
        ("Admin Login", "admin / AdminPassword123"),
        ("Official Login", "password same as configured username"),
    ]
    for label, value in cred_items:
        p = tf_links.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        run_label = p.add_run()
        run_label.text = f"{label}: "
        set_font(run_label, Pt(9), bold=True, color=DARK_TEXT)
        run_val = p.add_run()
        run_val.text = value
        set_font(run_val, Pt(9), color=GRAY_MED, name=FONT_MONO)

    # License
    tf4 = add_textbox(slide, Inches(0.5), Inches(4.9), Inches(9.0), Inches(0.4))
    tf4.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf4.paragraphs[0].add_run()
    run.text = d["license"]
    set_font(run, Pt(10), color=GRAY_MED, italic=True)

    return slide


# ============================================================
# MAIN
# ============================================================

def main():
    print(f"Loading template: {TEMPLATE_PATH}")
    prs = Presentation(TEMPLATE_PATH)

    slides = list(prs.slides)
    print(f"  Template has {len(slides)} slides")

    # Phase 1: Build content on template slides 1-13
    # Template slide 14 is LEFT UNTOUCHED (branded ending background)
    template_builders = [
        build_slide1,   # 1: Template Landing (fields only)
        build_slide2,   # 2: Problem Statement
        build_slide3,   # 3: Solution Overview
        build_slide4,   # 4: 5-Step Lifecycle
        build_slide5,   # 5: Citizen Portal
        build_slide6,   # 6: Official Portal
        build_slide7,   # 7: Admin Portal
        build_slide8,   # 8: Continuous Governance Loop
        build_slide9,   # 9: Why AI?
        build_slide10,  # 10: Voice Issue Reporting
        build_slide11,  # 11: AI Summarisation & Agenda
        build_slide12,  # 12: Face Authentication
        build_slide13,  # 13: Meeting Day Flow
    ]

    for i, builder in enumerate(template_builders):
        print(f"  Building slide {i + 1}: {builder.__doc__.strip()}")
        builder(slides[i])

    # Template slide 14 — DO NOT TOUCH (branded ending background)
    print(f"  Slide 14 (template ending): left untouched")

    # Phase 2: Add our eGramSabha title slide (extra)
    print(f"  Building extra: Our eGramSabha Title")
    build_our_title(prs)

    # Phase 3: Add extra content slides
    extra_builders = [
        (build_slide14, "TTS & Multilingual Support"),
        (build_slide15, "System Architecture - Diagram"),
        (build_slide16, "Architecture Layers"),
        (build_slide17, "Security Architecture"),
        (build_slide18, "Traffic Flow"),
        (build_slide19, "AWS Services"),
        (build_slide20, "Provider Abstraction"),
        (build_slide21, "Technologies Utilized"),
        (build_slide22, "Cost Breakdown"),
        (build_slide23, "Cost Optimizations"),
        (build_slide24, "Scale Projections"),
        (build_slide25, "Impact Potential"),
        (build_slide26, "Impact Outcomes"),
        (build_slide27, "National Impact"),
    ]

    for i, (builder, desc) in enumerate(extra_builders):
        print(f"  Building extra: {desc}")
        builder(prs)

    # Phase 3b: Add screenshot gallery slides (one per slide)
    print(f"  Building gallery: {len(GALLERY_SLIDES)} screenshot slides")
    build_gallery_slides(prs)
    for gs in GALLERY_SLIDES:
        print(f"    Gallery: {gs['title']}")

    # Phase 3c: Future Roadmap (after gallery)
    print(f"  Building extra: Future Roadmap")
    build_slide30(prs)

    # Phase 4: Add our Thank You / closing slide (extra)
    print(f"  Building extra: Our Thank You / Closing")
    build_closing_slide(prs)

    # Phase 5: XML slide reorder
    # Current order: [0..12]=template 1-13, [13]=template 14 (ending),
    #   [14]=our title, [15..31]=extra content, [32]=our closing
    # Target order:
    #   [0]=template landing, [1]=our title, [2..13]=template 2-13,
    #   [14..30]=extra content, [31]=our closing, [32]=template ending
    print(f"  Reordering slides...")
    slide_list = prs.slides._sldIdLst
    elements = list(slide_list)

    # Step 1: Remove template ending (index 13) and append to end
    ending_el = elements[13]
    slide_list.remove(ending_el)
    slide_list.append(ending_el)

    # Step 2: Our title is now at index 13 (first extra after template 1-13).
    # Move it to index 1 (right after template landing).
    elements = list(slide_list)
    title_el = elements[13]
    slide_list.remove(title_el)
    slide_list.insert(1, title_el)

    total = len(prs.slides)
    print(f"\n  Total slides: {total}")
    print(f"  Slide 1: Template landing (fields filled)")
    print(f"  Slide 2: Our eGramSabha title")
    print(f"  Slide {total - 1}: Our Thank You / closing")
    print(f"  Slide {total}: Template ending (untouched)")

    print(f"\nSaving to: {OUTPUT_PATH}")
    prs.save(OUTPUT_PATH)
    print("Done! PPTX generated successfully.")


if __name__ == "__main__":
    main()
