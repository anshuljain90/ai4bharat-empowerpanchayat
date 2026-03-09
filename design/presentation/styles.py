"""
Styling constants for eGramSabha Hackathon Presentation.
Based on template: Manrope font, #202729 text, 10"x5.625" slides.
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Slide dimensions (EMU)
SLIDE_WIDTH = 9144000   # 10 inches
SLIDE_HEIGHT = 5143500  # 5.625 inches

# Colors
DARK_TEXT = RGBColor(0x20, 0x27, 0x29)       # #202729 - primary text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_PRIMARY = RGBColor(0x2E, 0x7D, 0x32)   # Forest green (brand)
GREEN_LIGHT = RGBColor(0x4C, 0xAF, 0x50)     # Light green
ORANGE_ACCENT = RGBColor(0xFF, 0x6F, 0x00)   # Accent for highlights
BLUE_AWS = RGBColor(0x23, 0x2F, 0x3E)        # AWS dark blue
BLUE_LIGHT = RGBColor(0x00, 0x73, 0xBB)      # AWS service blue
GRAY_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)      # Table alternating
GRAY_MED = RGBColor(0x75, 0x75, 0x75)        # Subtitle gray
RED_URGENT = RGBColor(0xD3, 0x2F, 0x2F)      # Urgent/problem
YELLOW_READY = RGBColor(0xFF, 0xA0, 0x00)    # Code-ready services
TABLE_HEADER_BG = RGBColor(0x2E, 0x7D, 0x32) # Green header
TABLE_ALT_BG = RGBColor(0xE8, 0xF5, 0xE9)   # Light green alt rows
PLACEHOLDER_BG = RGBColor(0xE0, 0xE0, 0xE0)  # Screenshot placeholder
PLACEHOLDER_BORDER = RGBColor(0xBD, 0xBD, 0xBD)

# Extra slide constants
HEADER_BAR_COLOR = GREEN_PRIMARY
HEADER_BAR_HEIGHT = Inches(0.75)

# Font
FONT_TITLE = "Manrope"
FONT_BODY = "Manrope"
FONT_MONO = "Consolas"

# Standard font sizes
SIZE_SLIDE_TITLE = Pt(22)
SIZE_SUBTITLE = Pt(14)
SIZE_BODY = Pt(11)
SIZE_BODY_SMALL = Pt(9)
SIZE_TABLE = Pt(9)
SIZE_TABLE_HEADER = Pt(10)
SIZE_CAPTION = Pt(8)
SIZE_BIG_STAT = Pt(36)

# Content area (avoiding header banner area)
CONTENT_LEFT = Emu(311700)      # ~0.34"
CONTENT_TOP = Emu(863550)       # ~0.95" (below header)
CONTENT_WIDTH = Emu(8520600)    # ~9.33"
CONTENT_HEIGHT = Emu(3900000)   # ~4.27"

# Margins
MARGIN_LEFT = Inches(0.4)
MARGIN_TOP_TITLE = Inches(0.95)
MARGIN_TOP_BODY = Inches(1.5)
MARGIN_RIGHT = Inches(0.3)
BODY_WIDTH = Inches(9.2)
BODY_HEIGHT = Inches(3.8)


def set_font(run, size=SIZE_BODY, bold=False, color=DARK_TEXT, name=FONT_BODY, italic=False):
    """Set font properties on a run."""
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    run.font.italic = italic


def add_titled_textbox(slide, title_text, left, top, width, height,
                       title_size=SIZE_SLIDE_TITLE, title_color=DARK_TEXT):
    """Add a text box with a title paragraph, return the text frame for adding more."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    set_font(run, size=title_size, bold=True, color=title_color)
    return tf


def add_bullet_point(tf, text, level=0, size=SIZE_BODY, bold=False, color=DARK_TEXT):
    """Add a bullet point paragraph to a text frame."""
    p = tf.add_paragraph()
    p.level = level
    p.space_before = Pt(2)
    p.space_after = Pt(1)
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return p


def add_paragraph(tf, text, size=SIZE_BODY, bold=False, color=DARK_TEXT,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(2)):
    """Add a plain paragraph to a text frame."""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return p


def style_table(table, header_bg=TABLE_HEADER_BG, alt_bg=TABLE_ALT_BG):
    """Style a table with colored header and alternating rows."""
    for col_idx in range(len(table.columns)):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                set_font(run, size=SIZE_TABLE_HEADER, bold=True, color=WHITE)

    for row_idx in range(1, len(table.rows)):
        for col_idx in range(len(table.columns)):
            cell = table.cell(row_idx, col_idx)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_bg
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    set_font(run, size=SIZE_TABLE, color=DARK_TEXT)


def set_cell_text(table, row, col, text, size=SIZE_TABLE, bold=False, color=DARK_TEXT):
    """Set text in a table cell."""
    cell = table.cell(row, col)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.space_before = Pt(1)
    p.space_after = Pt(1)
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
